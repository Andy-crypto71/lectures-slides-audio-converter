import streamlit as st
import asyncio
import io
import os
import re
from pptx import Presentation
from docx import Document
from pypdf import PdfReader
import edge_tts

# --- INITIAL SETUP ---
st.set_page_config(
    page_title="Aura IA Studio - Study Audio Converter",
    page_icon="🎙️"
)

# --- LOGIN & SESSION LOGIC ---
def login_ui():
    st.markdown("<h1 style='text-align: center; color: #D4AF37;'>LysnUp</h1>", unsafe_allow_html=True)
    st.write("### Sign In to Continue")
    
    with st.container():
        email = st.text_input("Email")
        password = st.text_input("Password", type="password")
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("Login"):
                # For testing, we use "admin" and "aura123"
                if email == "admin@aura.com" and password == "aura123":
                    st.session_state['authenticated'] = True
                    st.rerun()
                else:
                    st.error("Invalid credentials")
        with col2:
            if st.button("Register"):
                st.info("Registration is currently limited to ambassadors.")

# --- TEXT CLEANING & EXTRACTION ---
def clean_text(text):
    text = re.sub(r'\[.*?\]', '', text)
    text = re.sub(r'\{.*?\}', '', text)
    symbols_to_remove = ['[', ']', '{', '}', '/', '\\', '_', '*', '#']
    for symbol in symbols_to_remove:
        text = text.replace(symbol, ' ')
    return re.sub(r'\s+', ' ', text).strip()

def extract_text(uploaded_file, ext):
    file_bytes = io.BytesIO(uploaded_file.read())
    content = []
    if ext == "pptx":
        prs = Presentation(file_bytes)
        for i, slide in enumerate(prs.slides):
            slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
            content.append(f"Slide {i + 1}. " + " ".join(slide_text))
    elif ext == "docx":
        doc = Document(file_bytes)
        content = [p.text for p in doc.paragraphs if p.text.strip()]
    elif ext == "pdf":
        reader = PdfReader(file_bytes)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text: content.append(f"Page {i + 1}. {text}")
    return " ... ".join(content)

async def save_audio(text, voice):
    communicate = edge_tts.Communicate(text, voice)
    await communicate.save("temp_output.mp3")
    return "temp_output.mp3"

# --- MAIN APP LOGIC ---
def main_app():
    # Logout button in sidebar
    st.sidebar.title("Aura IA Studio")
    if st.sidebar.button("Log Out"):
        st.session_state['authenticated'] = False
        st.rerun()

    st.title("🎙️ Study Audio Converter")
    st.write("Convert your lecture slides and notes into audio instantly.")

    voice = st.selectbox(
        "Select Voice",
        ["en-GB-ThomasNeural", "en-US-GuyNeural", "en-US-AriaNeural"]
    )

    file = st.file_uploader("Upload PPTX, DOCX, or PDF", type=["pptx", "docx", "pdf"])

    if file and st.button("Convert to Audio"):
        file_extension = file.name.split(".")[-1].lower()
        with st.spinner("Processing..."):
            raw_text = extract_text(file, file_extension)
            if raw_text:
                final_text = clean_text(raw_text)
                with st.spinner("Generating audio..."):
                    loop = asyncio.new_event_loop()
                    asyncio.set_event_loop(loop)
                    audio_file = loop.run_until_complete(save_audio(final_text, voice))
                    
                    st.success("Audio generated!")
                    st.audio(audio_file)

                    with open(audio_file, "rb") as f:
                        st.download_button(
                            label="Download MP3",
                            data=f,
                            file_name=f"{os.path.splitext(file.name)[0]}.mp3",
                            mime="audio/mpeg"
                        )
                    if os.path.exists(audio_file):
                        os.remove(audio_file)
            else:
                st.error("No readable text found.")

# --- ROUTING ---
if 'authenticated' not in st.session_state or not st.session_state['authenticated']:
    login_ui()
else:
    main_app()
